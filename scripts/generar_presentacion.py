"""Genera la presentacion del proyecto HHI como PDF (16:9) y PPTX.

Toda la composicion se hace en matplotlib (no PowerPoint templates). Cada
slide es una Figure de 16x9 pulgadas. Se exportan dos productos:

- artifacts/presentacion/presentacion_HHI.pdf  (un solo PDF multi-pagina)
- artifacts/presentacion/presentacion_HHI.pptx (PPTX con cada slide como
  imagen PNG en alta resolucion)

Datos: leidos de data/hhi_*.csv y constantes del informe (Bogota share, etc.).
"""

from __future__ import annotations

import io
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.backends.backend_pdf import PdfPages
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
DATA_DIR = ROOT / "data"
OUT_DIR = ROOT / "artifacts" / "presentacion"
TMP_PNG_DIR = OUT_DIR / "_slides_png"

# ---------------------------------------------------------------------------
# Estilo global
# ---------------------------------------------------------------------------

COLOR_PRIMARY = "#4E79A7"
COLOR_SECONDARY = "#76B7B2"
COLOR_WARNING = "#E0A800"
COLOR_ALERT = "#C43D3D"
COLOR_NEUTRAL = "#6B7280"
COLOR_TEXT = "#1F2937"
COLOR_TEXT_MUTED = "#4B5563"
COLOR_BG = "#FFFFFF"
COLOR_BG_SOFT = "#F7F8FA"
COLOR_HAIR = "#D1D5DB"

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "axes.titlesize": 26,
    "axes.labelsize": 20,
    "xtick.labelsize": 16,
    "ytick.labelsize": 16,
    "legend.fontsize": 16,
    "axes.edgecolor": COLOR_HAIR,
    "axes.linewidth": 0.6,
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.facecolor": COLOR_BG,
    "figure.facecolor": COLOR_BG,
    "savefig.facecolor": COLOR_BG,
})

SLIDE_W = 16
SLIDE_H = 9


# ---------------------------------------------------------------------------
# Datos
# ---------------------------------------------------------------------------

def load_data() -> dict[str, pd.DataFrame | dict]:
    anio = pd.read_csv(DATA_DIR / "hhi_por_anio.csv")
    nivel = pd.read_csv(DATA_DIR / "hhi_por_nivel.csv")
    master = pd.read_csv(DATA_DIR / "HHI_CRUCE_SECOP_DANE_RESULTADOS_final.csv")
    dept = pd.read_csv(DATA_DIR / "hhi_por_departamento.csv")

    # Cuota de Bogota en monto total (datos documentados en INFORME §5.5)
    bogota_share = pd.DataFrame({
        "anio": [2018, 2019, 2020, 2021, 2022, 2023, 2024],
        "cuota_pct": [50.2, 42.6, 41.2, 47.2, 54.3, 38.0, 34.1],
    })

    # KPIs del pipeline (verificados contra Tabla 2 del informe)
    pipeline_kpis = {
        "bronze_secop_i": 6_354_773,
        "bronze_secop_ii": 5_599_845,
        "bronze_cnpv_personas": 44_164_417,
        "silver_secop_i_tx": 5_456_438,
        "silver_secop_ii_tx": 4_026_650,
        "mart_filas": 13_860,
        "hhi_mercados": len(master),
        "anios_cubiertos": int(master["anio_key"].nunique()),
    }
    return {
        "anio": anio,
        "nivel": nivel,
        "master": master,
        "dept": dept,
        "bogota": bogota_share,
        "kpis": pipeline_kpis,
    }


# ---------------------------------------------------------------------------
# Componentes reutilizables
# ---------------------------------------------------------------------------

def base_slide() -> tuple[plt.Figure, plt.Axes]:
    fig = plt.figure(figsize=(SLIDE_W, SLIDE_H), dpi=150)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 9)
    ax.axis("off")
    # Barra decorativa superior
    ax.add_patch(mpatches.Rectangle((0, 8.78), 16, 0.22, color=COLOR_PRIMARY))
    return fig, ax


def slide_title(ax: plt.Axes, title: str, eyebrow: str = "") -> None:
    if eyebrow:
        ax.text(0.6, 7.95, eyebrow.upper(), fontsize=14,
                color=COLOR_PRIMARY, weight="bold")
    ax.text(0.6, 7.30, title, fontsize=42, color=COLOR_TEXT, weight="bold",
            va="top", wrap=True)


def slide_footer(ax: plt.Axes, slide_num: int, total: int) -> None:
    ax.text(0.6, 0.35,
            "Sinergia socioeconomica - HHI SECOP  ·  USTA 2026-I",
            fontsize=10, color=COLOR_NEUTRAL)
    ax.text(15.4, 0.35, f"{slide_num} / {total}",
            fontsize=10, color=COLOR_NEUTRAL, ha="right")


def add_chart_axes(fig: plt.Figure, rect: tuple[float, float, float, float]) -> plt.Axes:
    """Crea un ax para una grafica dentro del slide (coords en fracciones)."""
    return fig.add_axes(rect)


# ---------------------------------------------------------------------------
# Slides individuales
# ---------------------------------------------------------------------------

def slide_cover(data: dict, num: int, total: int) -> plt.Figure:
    fig = plt.figure(figsize=(SLIDE_W, SLIDE_H), dpi=150)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 9)
    ax.axis("off")
    # Bloque de color a la izquierda
    ax.add_patch(mpatches.Rectangle((0, 0), 5.5, 9, color=COLOR_PRIMARY))
    # Forma decorativa
    ax.add_patch(mpatches.Circle((5.5, 4.5), 2.2, color="#FFFFFF",
                                  alpha=0.08, zorder=2))
    ax.add_patch(mpatches.Circle((5.5, 4.5), 1.2, color="#FFFFFF",
                                  alpha=0.12, zorder=2))
    ax.text(0.5, 8.0, "USTA 2026-I", fontsize=18, color="#FFFFFF",
            weight="bold", alpha=0.85)
    ax.text(0.5, 1.2, "Consultorio de Estadistica\nObservatorio Ustadistica",
            fontsize=14, color="#FFFFFF", alpha=0.85)

    ax.text(6.2, 6.8, "Concentracion de la\ncontratacion publica\nen Colombia",
            fontsize=46, color=COLOR_TEXT, weight="bold", va="top",
            linespacing=1.1)
    ax.text(6.2, 3.5, "Indice Herfindahl-Hirschman aplicado a SECOP I + II",
            fontsize=22, color=COLOR_PRIMARY, weight="bold")
    ax.text(6.2, 2.9, "Periodo analitico 2018 - 2026",
            fontsize=18, color=COLOR_NEUTRAL)

    ax.text(6.2, 1.5, "Proyecto Sinergia socioeconomica",
            fontsize=14, color=COLOR_TEXT_MUTED, style="italic")
    ax.text(6.2, 1.1, "Modalidad: exposicion para dos presentadores",
            fontsize=12, color=COLOR_TEXT_MUTED)
    return fig


def slide_question(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax = base_slide()
    slide_title(ax, "La pregunta", eyebrow="Motivacion")
    ax.text(8, 4.7,
            "Cuando el Estado contrata,\n¿el valor que paga se reparte\nentre muchos proveedores\no se concentra en pocos?",
            fontsize=44, color=COLOR_TEXT, weight="bold",
            ha="center", va="center", linespacing=1.25)
    ax.text(8, 1.5,
            "Esa respuesta tiene implicaciones en competencia,\n"
            "eficiencia del gasto y riesgo de captura de rentas.",
            fontsize=20, color=COLOR_TEXT_MUTED, ha="center",
            va="center", linespacing=1.3)
    slide_footer(ax, num, total)
    return fig


def slide_fuentes(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax = base_slide()
    slide_title(ax, "Cinco fuentes de datos abiertos", eyebrow="Insumos")

    # Dos columnas
    ax.text(2.0, 6.2, "Protagonistas del HHI", fontsize=22,
            color=COLOR_PRIMARY, weight="bold")
    ax.text(2.0, 5.8, "Calculo del indicador", fontsize=13,
            color=COLOR_NEUTRAL, style="italic")
    proton = [
        ("SECOP I", "f789-7hwg", "Procesos Compra Publica", "Datos Abiertos Colombia"),
        ("SECOP II", "jbjy-vk9h", "Contratos Electronicos", "Datos Abiertos Colombia"),
    ]
    y = 4.7
    for nombre, ident, desc, src in proton:
        ax.add_patch(mpatches.FancyBboxPatch(
            (1.2, y - 0.95), 5.6, 1.15,
            boxstyle="round,pad=0.05,rounding_size=0.1",
            linewidth=0, facecolor=COLOR_PRIMARY, alpha=0.95))
        ax.text(1.5, y - 0.05, nombre, fontsize=24, color="#FFFFFF", weight="bold")
        ax.text(1.5, y - 0.55, desc, fontsize=14, color="#FFFFFF", alpha=0.9)
        ax.text(6.6, y - 0.05, ident, fontsize=14, color="#FFFFFF",
                alpha=0.9, ha="right", family="monospace")
        ax.text(6.6, y - 0.55, src, fontsize=11, color="#FFFFFF",
                alpha=0.75, ha="right")
        y -= 1.45

    ax.text(9.5, 6.2, "Contextuales", fontsize=22,
            color=COLOR_SECONDARY, weight="bold")
    ax.text(9.5, 5.8, "Enriquecen el cruce", fontsize=13,
            color=COLOR_NEUTRAL, style="italic")
    contextuales = [
        ("CNPV 2018", "Cat. 643", "Censo Poblacion y Vivienda", "DANE"),
        ("EMICRON", "Cat. 875", "Encuesta de Micronegocios", "DANE"),
        ("Proyecciones", "2018-2050", "Poblacion proyectada", "DANE"),
    ]
    y = 4.7
    for nombre, ident, desc, src in contextuales:
        ax.add_patch(mpatches.FancyBboxPatch(
            (8.9, y - 0.75), 6.4, 0.95,
            boxstyle="round,pad=0.05,rounding_size=0.1",
            linewidth=1.0, facecolor="#FFFFFF",
            edgecolor=COLOR_SECONDARY))
        ax.text(9.2, y - 0.05, nombre, fontsize=18, color=COLOR_TEXT, weight="bold")
        ax.text(9.2, y - 0.45, desc, fontsize=12, color=COLOR_NEUTRAL)
        ax.text(15.0, y - 0.05, ident, fontsize=12, color=COLOR_SECONDARY,
                ha="right", family="monospace", weight="bold")
        ax.text(15.0, y - 0.45, src, fontsize=10, color=COLOR_NEUTRAL, ha="right")
        y -= 1.15

    ax.text(8, 0.85,
            "El HHI se calcula con SECOP I + II. CNPV, EMICRON y proyecciones aportan contexto territorial.",
            fontsize=13, color=COLOR_TEXT_MUTED, ha="center", style="italic")
    slide_footer(ax, num, total)
    return fig


def slide_medallion(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax = base_slide()
    slide_title(ax, "Arquitectura Medallion", eyebrow="Pipeline")

    capas = [
        ("BRONZE", "Ingesta cruda\nParquet + metadatos\nHash de integridad",
         "#CD7F32", 2.0),
        ("SILVER", "Limpieza, tipos\nDIVIPOLA 5 digitos\nDeduplicacion",
         "#C0C0C0", 7.0),
        ("GOLD", "Modelo estrella\nFacts + dimensiones\nMart analitico",
         "#D4AF37", 12.0),
    ]
    box_w, box_h = 3.8, 3.4
    y_box = 2.5
    for i, (nombre, desc, color, x_center) in enumerate(capas):
        ax.add_patch(mpatches.FancyBboxPatch(
            (x_center - box_w / 2, y_box), box_w, box_h,
            boxstyle="round,pad=0.05,rounding_size=0.15",
            linewidth=0, facecolor=color, alpha=0.18))
        ax.add_patch(mpatches.Rectangle(
            (x_center - box_w / 2, y_box + box_h - 0.7), box_w, 0.7,
            facecolor=color, alpha=0.85))
        ax.text(x_center, y_box + box_h - 0.35, nombre, fontsize=24,
                color="#FFFFFF", weight="bold", ha="center", va="center")
        ax.text(x_center, y_box + 1.55, desc, fontsize=16,
                color=COLOR_TEXT, ha="center", va="center", linespacing=1.4)
        # Flecha entre capas
        if i < len(capas) - 1:
            ax.annotate("", xy=(capas[i + 1][3] - box_w / 2 - 0.2, y_box + box_h / 2),
                        xytext=(x_center + box_w / 2 + 0.2, y_box + box_h / 2),
                        arrowprops=dict(arrowstyle="->", color=COLOR_NEUTRAL,
                                        linewidth=2))

    ax.text(8, 1.2,
            "Trazabilidad: cualquier numero del HHI se puede rastrear hasta el contrato original.",
            fontsize=15, color=COLOR_TEXT_MUTED, ha="center", style="italic")
    slide_footer(ax, num, total)
    return fig


def slide_modelo_dim(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax = base_slide()
    slide_title(ax, "Modelo dimensional (esquema estrella)", eyebrow="Capa Gold")

    # Mart en el centro
    cx, cy = 8.0, 4.4
    ax.add_patch(mpatches.FancyBboxPatch(
        (cx - 1.8, cy - 0.8), 3.6, 1.6,
        boxstyle="round,pad=0.05,rounding_size=0.15",
        linewidth=0, facecolor=COLOR_PRIMARY))
    ax.text(cx, cy + 0.18, "MART", fontsize=22, color="#FFFFFF",
            weight="bold", ha="center", va="center")
    ax.text(cx, cy - 0.3, "13,860 filas", fontsize=14, color="#FFFFFF",
            ha="center", va="center", alpha=0.92)
    ax.text(cx, cy - 0.6, "municipio x anio", fontsize=11, color="#FFFFFF",
            ha="center", va="center", alpha=0.8, style="italic")

    # Dimensiones (arriba, bajadas para no chocar con titulo)
    dims = [
        ("dim_tiempo", "anios 2018-2029", 4.5, 6.0),
        ("dim_territorio", "1,155 DIVIPOLA", 11.5, 6.0),
    ]
    for nombre, desc, x, y in dims:
        ax.add_patch(mpatches.FancyBboxPatch(
            (x - 1.5, y - 0.6), 3.0, 1.2,
            boxstyle="round,pad=0.05,rounding_size=0.1",
            linewidth=1.4, facecolor="#FFFFFF", edgecolor=COLOR_SECONDARY))
        ax.text(x, y + 0.18, nombre, fontsize=16, color=COLOR_SECONDARY,
                weight="bold", ha="center", va="center", family="monospace")
        ax.text(x, y - 0.25, desc, fontsize=12, color=COLOR_NEUTRAL,
                ha="center", va="center")
        ax.plot([x, cx], [y - 0.6, cy + 0.8], color=COLOR_HAIR, linewidth=1)

    # Facts (abajo)
    facts = [
        ("fact_contratacion", "SECOP I + II", 2.5, 1.7),
        ("fact_censo", "CNPV 2018", 6.5, 1.5),
        ("fact_micronegocios", "EMICRON", 9.5, 1.5),
        ("fact_demografia", "Proyecciones", 13.5, 1.7),
    ]
    for nombre, src, x, y in facts:
        color = COLOR_ALERT if "contratacion" in nombre else COLOR_NEUTRAL
        ax.add_patch(mpatches.FancyBboxPatch(
            (x - 1.4, y - 0.5), 2.8, 1.0,
            boxstyle="round,pad=0.05,rounding_size=0.1",
            linewidth=1.4, facecolor="#FFFFFF", edgecolor=color))
        ax.text(x, y + 0.15, nombre, fontsize=13,
                color=color, weight="bold", ha="center", va="center",
                family="monospace")
        ax.text(x, y - 0.25, src, fontsize=11, color=COLOR_NEUTRAL,
                ha="center", va="center")
        ax.plot([x, cx], [y + 0.5, cy - 0.8], color=COLOR_HAIR, linewidth=1)

    ax.text(2.5, 0.85, "★ La fuente del HHI",
            fontsize=12, color=COLOR_ALERT, ha="center", style="italic",
            weight="bold")
    slide_footer(ax, num, total)
    return fig


def slide_pipeline_kpis(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax = base_slide()
    slide_title(ax, "Del CSV crudo al indicador", eyebrow="Volumen procesado")

    kpis = [
        ("11.95 M", "filas SECOP I + II\ningestadas en Bronze", COLOR_PRIMARY),
        ("9.48 M", "contratos limpios\nen Silver transaccional", COLOR_SECONDARY),
        ("13,860", "filas del mart\nGold (municipio x anio)", COLOR_WARNING),
        ("11,792", "mercados unicos\ncon HHI calculado", COLOR_ALERT),
    ]
    x_positions = [2.0, 6.0, 10.0, 14.0]
    for (kpi, label, color), x in zip(kpis, x_positions):
        ax.add_patch(mpatches.Circle((x, 5.3), 1.45, color=color, alpha=0.15))
        ax.add_patch(mpatches.Circle((x, 5.3), 1.2, color=color, alpha=0.25))
        ax.text(x, 5.3, kpi, fontsize=34, color=color, weight="bold",
                ha="center", va="center")
        ax.text(x, 3.2, label, fontsize=15, color=COLOR_TEXT,
                ha="center", va="center", linespacing=1.3)

    ax.annotate("", xy=(5.0, 5.3), xytext=(3.45, 5.3),
                arrowprops=dict(arrowstyle="->", color=COLOR_NEUTRAL, linewidth=2))
    ax.annotate("", xy=(9.0, 5.3), xytext=(7.45, 5.3),
                arrowprops=dict(arrowstyle="->", color=COLOR_NEUTRAL, linewidth=2))
    ax.annotate("", xy=(13.0, 5.3), xytext=(11.45, 5.3),
                arrowprops=dict(arrowstyle="->", color=COLOR_NEUTRAL, linewidth=2))

    ax.text(8, 1.2,
            "Cobertura: 9 anos (2018-2026), 1,093 municipios con contratacion observada.",
            fontsize=15, color=COLOR_TEXT_MUTED, ha="center", style="italic")
    slide_footer(ax, num, total)
    return fig


def slide_que_es_hhi(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax = base_slide()
    slide_title(ax, "Que es el HHI", eyebrow="Definicion")

    # Formula a la izquierda
    ax.text(4.0, 5.5,
            r"$HHI_m = \sum_{i=1}^{n_m} s_{im}^{\,2}$",
            fontsize=48, color=COLOR_TEXT, ha="center", va="center")
    ax.text(4.0, 4.3,
            r"$s_{im}$ = participacion % del proveedor $i$",
            fontsize=18, color=COLOR_NEUTRAL, ha="center", va="center")
    ax.text(4.0, 3.7,
            "Mercado m = municipio x anio x orden de entidad",
            fontsize=15, color=COLOR_NEUTRAL, ha="center", va="center",
            style="italic")

    # Escala a la derecha
    bandas = [
        (0, 1500, "Baja\nconcentracion", COLOR_PRIMARY),
        (1500, 2500, "Concentracion\nmoderada", COLOR_WARNING),
        (2500, 10000, "Alta\nconcentracion", COLOR_ALERT),
    ]
    bar_x = 9.0
    bar_w = 5.5
    bar_h = 0.7
    bar_y = 5.2
    for lo, hi, label, color in bandas:
        frac_lo = lo / 10000
        frac_hi = hi / 10000
        x0 = bar_x + frac_lo * bar_w
        w = (frac_hi - frac_lo) * bar_w
        ax.add_patch(mpatches.Rectangle((x0, bar_y), w, bar_h,
                                         color=color, alpha=0.85))
        ax.text(x0 + w / 2, bar_y + bar_h / 2, label, fontsize=12,
                color="#FFFFFF", ha="center", va="center", weight="bold",
                linespacing=1.1)

    ax.text(bar_x, bar_y + bar_h + 0.3, "0", fontsize=12, color=COLOR_NEUTRAL)
    ax.text(bar_x + bar_w * 0.15, bar_y + bar_h + 0.3, "1,500",
            fontsize=12, color=COLOR_NEUTRAL, ha="center")
    ax.text(bar_x + bar_w * 0.25, bar_y + bar_h + 0.3, "2,500",
            fontsize=12, color=COLOR_NEUTRAL, ha="center")
    ax.text(bar_x + bar_w, bar_y + bar_h + 0.3, "10,000",
            fontsize=12, color=COLOR_NEUTRAL, ha="right")

    ax.text(bar_x + bar_w / 2, bar_y - 0.8,
            "Escala estandar (DOJ / FTC, Horizontal Merger Guidelines)",
            fontsize=13, color=COLOR_TEXT_MUTED, ha="center", style="italic")

    ax.text(bar_x + bar_w / 2, 3.3,
            "10,000 = un solo proveedor concentra el 100% del valor",
            fontsize=14, color=COLOR_ALERT, ha="center", weight="bold")

    slide_footer(ax, num, total)
    return fig


def slide_tendencia_anual(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax_bg = base_slide()
    slide_title(ax_bg, "Tendencia anual del HHI promedio", eyebrow="Resultado 1")
    slide_footer(ax_bg, num, total)
    ax_bg.text(0.6, 6.45,
               "La concentracion promedio nacional se mantiene en la banda BAJA-MODERADA.",
               fontsize=17, color=COLOR_PRIMARY, weight="bold")

    ax = add_chart_axes(fig, [0.08, 0.12, 0.84, 0.55])
    anio = data["anio"].sort_values("anio_key")
    x = anio["anio_key"].values
    y_avg = anio["HHI_promedio"].values
    y_med = anio["HHI_mediana"].values

    ax.fill_between(x, 0, y_avg, color=COLOR_PRIMARY, alpha=0.10)
    ax.plot(x, y_avg, marker="o", color=COLOR_PRIMARY, linewidth=3.5,
            markersize=11, label="HHI promedio")
    ax.plot(x, y_med, marker="s", color=COLOR_SECONDARY, linewidth=2.5,
            markersize=9, linestyle="--", label="HHI mediana")

    # Marcas
    ax.axhline(1500, color=COLOR_WARNING, linestyle=":", linewidth=1.4, alpha=0.7)
    ax.axhline(2500, color=COLOR_ALERT, linestyle=":", linewidth=1.4, alpha=0.7)
    ax.text(x[-1] + 0.1, 1500, " moderada", color=COLOR_WARNING, fontsize=12,
            va="center")
    ax.text(x[-1] + 0.1, 2500, " alta", color=COLOR_ALERT, fontsize=12,
            va="center")

    # Etiquetas en cada punto
    for xi, yi in zip(x, y_avg):
        ax.annotate(f"{yi:,.0f}", (xi, yi), textcoords="offset points",
                    xytext=(0, 14), ha="center", fontsize=13,
                    color=COLOR_TEXT, weight="bold")

    ax.set_xlabel("Anio")
    ax.set_ylabel("HHI")
    ax.set_ylim(0, max(y_avg.max() * 1.25, 2700))
    ax.legend(loc="upper right", frameon=False, fontsize=14)
    ax.grid(alpha=0.2)
    ax.set_xticks(x)
    return fig


def slide_nacional_vs_terr(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax_bg = base_slide()
    slide_title(ax_bg, "Nacional vs Territorial", eyebrow="Resultado 2")
    slide_footer(ax_bg, num, total)
    ax_bg.text(0.6, 6.45,
               "El orden NACIONAL concentra de forma sistematica MAS que el TERRITORIAL.",
               fontsize=17, color=COLOR_ALERT, weight="bold")

    nivel = data["nivel"]
    pivot = nivel.pivot_table(index="anio_key", columns="orden_entidad",
                              values="HHI_promedio", aggfunc="mean")

    ax = add_chart_axes(fig, [0.08, 0.12, 0.84, 0.55])
    if "NACIONAL" in pivot.columns:
        ax.plot(pivot.index, pivot["NACIONAL"], marker="o", color=COLOR_ALERT,
                linewidth=4, markersize=12, label="NACIONAL")
    if "TERRITORIAL" in pivot.columns:
        ax.plot(pivot.index, pivot["TERRITORIAL"], marker="o", color=COLOR_PRIMARY,
                linewidth=4, markersize=12, label="TERRITORIAL")

    ax.axhline(1500, color=COLOR_WARNING, linestyle=":", linewidth=1.4, alpha=0.7)
    ax.axhline(2500, color=COLOR_ALERT, linestyle=":", linewidth=1.4, alpha=0.7)
    ax.text(pivot.index[-1] + 0.1, 1500, " moderada", color=COLOR_WARNING,
            fontsize=12, va="center")
    ax.text(pivot.index[-1] + 0.1, 2500, " alta", color=COLOR_ALERT,
            fontsize=12, va="center")

    if "NACIONAL" in pivot.columns:
        last = pivot["NACIONAL"].dropna().iloc[-1]
        ax.annotate(f"{last:,.0f}",
                    (pivot["NACIONAL"].dropna().index[-1], last),
                    textcoords="offset points", xytext=(8, 0),
                    fontsize=15, color=COLOR_ALERT, weight="bold")
    if "TERRITORIAL" in pivot.columns:
        last_t = pivot["TERRITORIAL"].dropna().iloc[-1]
        ax.annotate(f"{last_t:,.0f}",
                    (pivot["TERRITORIAL"].dropna().index[-1], last_t),
                    textcoords="offset points", xytext=(8, 0),
                    fontsize=15, color=COLOR_PRIMARY, weight="bold")

    ax.set_xlabel("Anio")
    ax.set_ylabel("HHI promedio")
    ax.legend(loc="upper left", frameon=False, fontsize=16)
    ax.grid(alpha=0.2)
    ax.set_xticks(pivot.index)
    ax.set_xticklabels([str(int(i)) for i in pivot.index])
    return fig


def slide_distribucion(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax_bg = base_slide()
    slide_title(ax_bg, "Distribucion del HHI", eyebrow="Resultado 3")
    slide_footer(ax_bg, num, total)
    ax_bg.text(0.6, 6.45,
               "La mayoria de mercados (mediana 644) esta en la banda BAJA. La cola alta es minoritaria.",
               fontsize=17, color=COLOR_PRIMARY, weight="bold")

    master = data["master"]
    ax = add_chart_axes(fig, [0.08, 0.12, 0.84, 0.55])
    bins = np.linspace(0, 10000, 41)
    counts, edges, patches = ax.hist(master["HHI"], bins=bins,
                                     color=COLOR_PRIMARY, edgecolor="white",
                                     alpha=0.9)
    # Pinta el bin maximo (HHI=10000) en rojo
    for patch, edge_left in zip(patches, edges[:-1]):
        if edge_left >= 9750:
            patch.set_facecolor(COLOR_ALERT)

    ax.axvline(1500, color=COLOR_WARNING, linestyle="--", linewidth=1.6)
    ax.axvline(2500, color=COLOR_ALERT, linestyle="--", linewidth=1.6)
    ax.text(1500, ax.get_ylim()[1] * 0.95, "  1,500", color=COLOR_WARNING,
            fontsize=13, va="top")
    ax.text(2500, ax.get_ylim()[1] * 0.95, "  2,500", color=COLOR_ALERT,
            fontsize=13, va="top")

    n_10000 = (master["HHI"] == 10000).sum()
    ax.annotate(f"HHI = 10,000\n{n_10000} mercados ({n_10000/len(master)*100:.2f}%)",
                xy=(9900, 5), xytext=(8000, ax.get_ylim()[1] * 0.55),
                fontsize=14, color=COLOR_ALERT, weight="bold",
                ha="center", linespacing=1.3,
                arrowprops=dict(arrowstyle="->", color=COLOR_ALERT, linewidth=1.6))

    ax.set_xlabel("HHI")
    ax.set_ylabel("Mercados")
    ax.grid(axis="y", alpha=0.2)
    return fig


def slide_top_dept(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax_bg = base_slide()
    slide_title(ax_bg, "Top departamentos con mayor HHI (2026)",
                eyebrow="Resultado 4")
    slide_footer(ax_bg, num, total)
    ax_bg.text(0.6, 6.45,
               "Cinco departamentos cruzan la banda ALTA (>2,500): Atlantico y Choco al frente.",
               fontsize=17, color=COLOR_ALERT, weight="bold")

    master = data["master"]
    ultimo = master["anio_key"].max()
    top = (master[master["anio_key"] == ultimo]
           .groupby("nombre_departamento", dropna=True)["HHI"]
           .mean().sort_values(ascending=True).tail(10))

    ax = add_chart_axes(fig, [0.20, 0.12, 0.72, 0.55])
    colors = []
    for v in top.values:
        if v >= 2500:
            colors.append(COLOR_ALERT)
        elif v >= 1500:
            colors.append(COLOR_WARNING)
        else:
            colors.append(COLOR_PRIMARY)
    bars = ax.barh(top.index, top.values, color=colors, edgecolor="white", height=0.75)
    for bar, v in zip(bars, top.values):
        ax.text(v + 50, bar.get_y() + bar.get_height() / 2,
                f"{v:,.0f}", va="center", fontsize=14,
                color=COLOR_TEXT, weight="bold")

    ax.axvline(1500, color=COLOR_WARNING, linestyle="--", linewidth=1, alpha=0.6)
    ax.axvline(2500, color=COLOR_ALERT, linestyle="--", linewidth=1, alpha=0.6)
    ax.set_xlabel("HHI promedio")
    ax.set_xlim(0, max(top.values) * 1.18)
    ax.tick_params(axis="y", labelsize=15)
    ax.grid(axis="x", alpha=0.2)
    return fig


def slide_186(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax = base_slide()
    slide_title(ax, "Los 186 mercados con HHI = 10,000", eyebrow="Aclaracion")

    master = data["master"]
    n_total = len(master)
    n_10000 = int((master["HHI"] == 10000).sum())
    trivial = int(((master["HHI"] == 10000) & (master["total_contratos"] == 1)).sum())
    real = n_10000 - trivial

    ax.text(8, 6.5, f"{n_10000} mercados", fontsize=64, color=COLOR_ALERT,
            ha="center", va="center", weight="bold")
    ax.text(8, 5.5,
            f"de {n_total:,} totales  =  {n_10000/n_total*100:.2f}%  (menos del 2%)",
            fontsize=22, color=COLOR_TEXT, ha="center", va="center")

    # Tres cards
    cards = [
        (trivial, "1 contrato + 1 proveedor",
         "Resultado matematicamente trivial.\nNo es informativo.", COLOR_NEUTRAL),
        (real, ">= 2 contratos al mismo NIT",
         "Monopolios REALES.\nCasos para revision cualitativa.", COLOR_ALERT),
        (0, "HHI=10,000 con > 1 proveedor",
         "Imposible bajo la formula.\nValidacion del calculo OK.", COLOR_SECONDARY),
    ]
    x_starts = [0.6, 6.0, 11.4]
    for (n, t, d, color), x in zip(cards, x_starts):
        ax.add_patch(mpatches.FancyBboxPatch(
            (x, 1.0), 4.0, 3.5,
            boxstyle="round,pad=0.05,rounding_size=0.15",
            linewidth=0, facecolor=color, alpha=0.12))
        ax.add_patch(mpatches.Rectangle(
            (x, 4.30), 4.0, 0.20, facecolor=color))
        ax.text(x + 2.0, 3.85, f"{n}", fontsize=42, color=color,
                weight="bold", ha="center", va="center")
        ax.text(x + 2.0, 2.85, t, fontsize=15, color=COLOR_TEXT,
                weight="bold", ha="center", va="center")
        ax.text(x + 2.0, 1.8, d, fontsize=12, color=COLOR_TEXT_MUTED,
                ha="center", va="center", linespacing=1.4)

    slide_footer(ax, num, total)
    return fig


def slide_sesgo_bogota(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax_bg = base_slide()
    slide_title(ax_bg, "Sesgo geografico: Bogota concentra",
                eyebrow="Limitacion critica")
    slide_footer(ax_bg, num, total)
    ax_bg.text(0.6, 6.45,
               "DIVIPOLA en SECOP = municipio de la ENTIDAD, no del lugar de ejecucion.",
               fontsize=17, color=COLOR_ALERT, weight="bold")
    ax_bg.text(0.6, 5.95,
               "Las entidades del orden nacional tienen sede en Bogota. Mitigacion: segmentar por orden_entidad.",
               fontsize=13, color=COLOR_TEXT_MUTED, style="italic")

    bog = data["bogota"]
    ax = add_chart_axes(fig, [0.10, 0.12, 0.82, 0.50])
    bars = ax.bar(bog["anio"], bog["cuota_pct"], color=COLOR_ALERT,
                  edgecolor="white", width=0.65)
    for bar, v in zip(bars, bog["cuota_pct"]):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 1.5,
                f"{v:.1f}%", ha="center", fontsize=14,
                color=COLOR_ALERT, weight="bold")

    ax.set_ylim(0, 65)
    ax.set_ylabel("Cuota Bogota (%)")
    ax.set_xlabel("Anio")
    ax.set_xticks(bog["anio"])
    ax.grid(axis="y", alpha=0.2)
    ax.axhline(50, color=COLOR_NEUTRAL, linestyle=":", linewidth=1, alpha=0.6)
    return fig


def slide_hallazgos(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax = base_slide()
    slide_title(ax, "Cuatro hallazgos clave", eyebrow="Sintesis")

    hallazgos = [
        ("1", "Baja-moderada",
         "El HHI promedio nacional se mantiene entre 1,040 y 1,484 en 9 anos.",
         COLOR_PRIMARY),
        ("2", "Nacional > Territorial",
         "El orden nacional duplica o triplica el HHI del territorial en todos los anos.",
         COLOR_ALERT),
        ("3", "Focos territoriales",
         "Atlantico, Choco, Magdalena, La Guajira y Boyaca cruzan la banda alta en 2026.",
         COLOR_WARNING),
        ("4", "Solo 1.58% al maximo",
         "186 mercados con HHI=10,000; de esos solo 19 son monopolios estructurales.",
         COLOR_SECONDARY),
    ]
    positions = [(0.6, 4.5), (8.2, 4.5), (0.6, 0.7), (8.2, 0.7)]
    for (idx, tit, desc, color), (x, y) in zip(hallazgos, positions):
        ax.add_patch(mpatches.FancyBboxPatch(
            (x, y), 7.2, 3.5,
            boxstyle="round,pad=0.05,rounding_size=0.15",
            linewidth=0, facecolor=color, alpha=0.10))
        ax.add_patch(mpatches.Circle((x + 0.85, y + 2.6), 0.55, color=color))
        ax.text(x + 0.85, y + 2.6, idx, fontsize=28, color="#FFFFFF",
                weight="bold", ha="center", va="center")
        ax.text(x + 1.7, y + 2.7, tit, fontsize=22, color=color,
                weight="bold", va="center")
        ax.text(x + 0.3, y + 1.45, desc, fontsize=15, color=COLOR_TEXT,
                wrap=True, linespacing=1.5)

    slide_footer(ax, num, total)
    return fig


def slide_limitaciones(data: dict, num: int, total: int) -> plt.Figure:
    fig, ax = base_slide()
    slide_title(ax, "Limitaciones explicitas", eyebrow="Honestidad estadistica")

    lims = [
        ("Sesgo geografico de SECOP",
         "La DIVIPOLA representa la entidad contratante,\nno el lugar de ejecucion del contrato."),
        ("Concentracion del valor, no de competencia",
         "El HHI no mide numero de oferentes\nni pluralidad real de propuestas."),
        ("Restriccion legal DANE",
         "Secreto estadistico (Ley 79/1993):\nno se cruza NIT con personas del censo."),
    ]
    x_starts = [0.6, 6.0, 11.4]
    for (tit, desc), x in zip(lims, x_starts):
        ax.add_patch(mpatches.FancyBboxPatch(
            (x, 2.5), 4.0, 4.2,
            boxstyle="round,pad=0.05,rounding_size=0.15",
            linewidth=2, facecolor="#FFFFFF", edgecolor=COLOR_WARNING))
        ax.add_patch(mpatches.Rectangle((x, 6.50), 4.0, 0.2,
                                         facecolor=COLOR_WARNING))
        ax.text(x + 2.0, 5.85, tit, fontsize=17, color=COLOR_TEXT,
                weight="bold", ha="center", va="center", wrap=True)
        ax.text(x + 2.0, 3.9, desc, fontsize=14, color=COLOR_TEXT_MUTED,
                ha="center", va="center", linespacing=1.5)

    ax.text(8, 1.4,
            "Las limitaciones se reportan SIEMPRE junto a los resultados.\nUn indicador honesto requiere caveats explicitos.",
            fontsize=14, color=COLOR_TEXT_MUTED, ha="center",
            style="italic", linespacing=1.4)
    slide_footer(ax, num, total)
    return fig


def slide_cierre(data: dict, num: int, total: int) -> plt.Figure:
    fig = plt.figure(figsize=(SLIDE_W, SLIDE_H), dpi=150)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 9)
    ax.axis("off")
    ax.add_patch(mpatches.Rectangle((0, 0), 16, 9, color=COLOR_PRIMARY))
    ax.add_patch(mpatches.Circle((12, 7), 3, color="#FFFFFF", alpha=0.06))
    ax.add_patch(mpatches.Circle((4, 2), 4, color="#FFFFFF", alpha=0.05))

    ax.text(8, 5.8, "Gracias", fontsize=88, color="#FFFFFF",
            ha="center", va="center", weight="bold")
    ax.text(8, 4.3, "Preguntas y comentarios", fontsize=24,
            color="#FFFFFF", ha="center", va="center", alpha=0.9)

    ax.text(8, 2.5, "Sinergia socioeconomica - Concentracion de la contratacion publica",
            fontsize=15, color="#FFFFFF", ha="center", alpha=0.85)
    ax.text(8, 2.0, "Consultorio de Estadistica USTA - Observatorio Ustadistica 2026-I",
            fontsize=13, color="#FFFFFF", ha="center", alpha=0.75)
    ax.text(8, 1.0,
            "Repositorio: desarrollo_social_y_economico   ·   Informe: docs/INFORME_HHI_DETALLADO.md",
            fontsize=11, color="#FFFFFF", ha="center", alpha=0.65,
            family="monospace")

    return fig


# ---------------------------------------------------------------------------
# Render
# ---------------------------------------------------------------------------

SLIDE_FUNCTIONS = [
    slide_cover,
    slide_question,
    slide_fuentes,
    slide_medallion,
    slide_modelo_dim,
    slide_pipeline_kpis,
    slide_que_es_hhi,
    slide_tendencia_anual,
    slide_nacional_vs_terr,
    slide_distribucion,
    slide_top_dept,
    slide_186,
    slide_sesgo_bogota,
    slide_hallazgos,
    slide_limitaciones,
    slide_cierre,
]


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    TMP_PNG_DIR.mkdir(parents=True, exist_ok=True)

    data = load_data()
    total = len(SLIDE_FUNCTIONS)
    pdf_path = OUT_DIR / "presentacion_HHI.pdf"

    # PDF multi-pagina
    with PdfPages(pdf_path) as pdf:
        for i, fn in enumerate(SLIDE_FUNCTIONS, start=1):
            fig = fn(data, i, total)
            pdf.savefig(fig, dpi=150)
            # Tambien lo guardo como PNG para el PPTX
            png_path = TMP_PNG_DIR / f"slide_{i:02d}.png"
            fig.savefig(png_path, dpi=170, facecolor=fig.get_facecolor())
            plt.close(fig)

    print(f"OK PDF -> {pdf_path.relative_to(ROOT)}")

    # PPTX con cada slide como imagen
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]

    for i in range(1, total + 1):
        slide = prs.slides.add_slide(blank_layout)
        png = TMP_PNG_DIR / f"slide_{i:02d}.png"
        slide.shapes.add_picture(str(png), 0, 0,
                                  width=prs.slide_width,
                                  height=prs.slide_height)

    pptx_path = OUT_DIR / "presentacion_HHI.pptx"
    prs.save(str(pptx_path))
    print(f"OK PPTX -> {pptx_path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
